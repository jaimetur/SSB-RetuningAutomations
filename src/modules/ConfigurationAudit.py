# -*- coding: utf-8 -*-

import os
import re
from typing import List, Tuple, Optional, Dict
import pandas as pd

from src.utils.Utils import get_resource_path

from src.modules.CommonMethods import (
    read_text_with_encoding,
    find_all_subnetwork_headers,
    extract_mo_from_subnetwork_line,
    parse_table_slice_from_subnetwork,
    SUMMARY_RE,
    sanitize_sheet_name,
    unique_sheet_name,
    natural_logfile_key,
    color_summary_tabs,
    enable_header_filters,
)


class ConfigurationAudit:
    """
    Generates an Excel in input_dir with one sheet per *.log / *.logs / *.txt file.
    (Functionality kept, extended with SummaryAudit sheet and PPT summary.)

    ARFCN-related parameters (N77, etc.) are now configurable via __init__:
      - new_arfcn            → main "new" NR / LTE ARFCN (e.g. 648672)
      - old_arfcn            → main "old" NR / LTE ARFCN (e.g. 647328)
      - allowed_n77_ssb      → allowed SSB values for N77 (e.g. {648672, 653952})
      - allowed_n77_arfcn    → allowed ARFCN values for N77 sectors
    """

    SUMMARY_RE = SUMMARY_RE  # keep class reference

    def __init__(
        self,
        old_arfcn: int,
        new_arfcn: int,
        allowed_n77_ssb: Optional[List[int]] = None,
        allowed_n77_arfcn: Optional[List[int]] = None,
    ):
        """
        Initialize ConfigurationAudit with ARFCN-related parameters.

        All values are converted to integers/sets of integers internally to make checks robust.
        """
        # Core ARFCN values
        self.OLD_ARFCN: int = int(old_arfcn)
        self.NEW_ARFCN: int = int(new_arfcn)

        # Allowed SSB values for N77 cells (e.g. {648672, 653952})
        if allowed_n77_ssb is None:
            self.ALLOWED_N77_SSB = set()
        else:
            self.ALLOWED_N77_SSB = {int(v) for v in allowed_n77_ssb}

        # Allowed ARFCN values for N77B sectors (e.g. {654652, 655324, 655984, 656656})
        if allowed_n77_arfcn is None:
            self.ALLOWED_N77_ARFCN = set()
        else:
            self.ALLOWED_N77_ARFCN = {int(v) for v in allowed_n77_arfcn}

    # =====================================================================
    #                            PUBLIC API
    # =====================================================================
    def run(
        self,
        input_dir: str,
        module_name: Optional[str] = "",
        versioned_suffix: Optional[str] = None,
        tables_order: Optional[List[str]] = None,      # optional sheet ordering
        filter_frequencies: Optional[List[str]] = None # substrings to filter pivot columns
    ) -> str:
        """
        Main entry point: creates an Excel file with one sheet per detected table.
        Sheets are ordered according to TABLES_ORDER if provided; otherwise,
        they are sorted in a natural order by filename (Data_Collection.txt, Data_Collection(1).txt, ...).

        If 'filter_frequencies' is provided, the three added summary sheets will keep only
        those pivot *columns* whose header contains any of the provided substrings
        (case-insensitive). 'NodeId' and 'Total' are always kept.

        In addition, a 'SummaryAudit' sheet is created with high-level checks
        across the parsed tables, and a PowerPoint (.pptx) summary is generated
        with a textual bullet-style overview per category.
        """
        # --- Normalize filters ---
        freq_filters = [str(f).strip() for f in (filter_frequencies or []) if str(f).strip()]

        # --- Validate the input directory ---
        if not os.path.isdir(input_dir):
            raise NotADirectoryError(f"Invalid directory: {input_dir}")

        # --- Detect log/txt files ---
        log_files = self._find_log_files(input_dir)
        if not log_files:
            raise FileNotFoundError(f"No .log/.logs/.txt files found in: {input_dir}")

        # --- Natural sorting of files (handles '(1)', '(2)', '(10)', etc.) ---
        sorted_files = sorted(log_files, key=natural_logfile_key)
        file_rank: Dict[str, int] = {os.path.basename(p): i for i, p in enumerate(sorted_files)}

        # --- Build MO (table) ranking if TABLES_ORDER is provided ---
        mo_rank: Dict[str, int] = {}
        if tables_order:
            mo_rank = {name: i for i, name in enumerate(tables_order)}

        # --- Prepare Excel output path ---
        excel_path = os.path.join(input_dir, f"ConfigurationAudit{versioned_suffix}.xlsx")
        table_entries: List[Dict[str, object]] = []

        # --- Keep a per-file index to preserve order of multiple tables inside same file ---
        per_file_table_idx: Dict[str, int] = {}

        # =====================================================================
        #                PHASE 1: Parse all log/txt files
        # =====================================================================
        for path in log_files:
            base_filename = os.path.basename(path)
            lines, encoding_used = self._read_text_file(path)

            header_indices = self._find_all_subnetwork_headers(lines)

            # --- Case 1: No 'SubNetwork' header found, fallback single-table mode ---
            if not header_indices:
                header_idx = self._find_subnetwork_header_index(lines)
                mo_name_prev = self._extract_mo_name_from_previous_line(lines, header_idx)
                df, note = self._parse_log_lines(lines, forced_header_idx=header_idx)

                if encoding_used:
                    note = (note + " | " if note else "") + f"encoding={encoding_used}"
                df, note = self._cap_rows(df, note)

                idx_in_file = per_file_table_idx.get(base_filename, 0)
                per_file_table_idx[base_filename] = idx_in_file + 1

                table_entries.append({
                    "df": df,
                    "sheet_candidate": mo_name_prev if mo_name_prev else os.path.splitext(base_filename)[0],
                    "log_file": base_filename,
                    "tables_in_log": 1,
                    "note": note or "",
                    "idx_in_file": idx_in_file,  # numeric index of this table inside the same file
                })
                continue

            # --- Case 2: Multiple 'SubNetwork' headers found (multi-table log) ---
            tables_in_log = len(header_indices)
            header_indices.append(len(lines))  # add sentinel index

            for ix in range(tables_in_log):
                h = header_indices[ix]
                nxt = header_indices[ix + 1]
                mo_name_from_line = extract_mo_from_subnetwork_line(lines[h])
                desired_sheet = mo_name_from_line if mo_name_from_line else os.path.splitext(base_filename)[0]

                df = parse_table_slice_from_subnetwork(lines, h, nxt)
                note = "Slice parsed"
                if encoding_used:
                    note += f" | encoding={encoding_used}"
                df, note = self._cap_rows(df, note)

                idx_in_file = per_file_table_idx.get(base_filename, 0)
                per_file_table_idx[base_filename] = idx_in_file + 1

                table_entries.append({
                    "df": df,
                    "sheet_candidate": desired_sheet,
                    "log_file": base_filename,
                    "tables_in_log": tables_in_log,
                    "note": note or "",
                    "idx_in_file": idx_in_file,
                })

        # =====================================================================
        #                PHASE 2: Determine final sorting order
        # =====================================================================
        def entry_sort_key(entry: Dict[str, object]) -> Tuple[int, int, int]:
            """
            Final sorting key for Excel sheets:
              - If TABLES_ORDER exists → sort by table order first, then by file (natural), then by table index
              - Otherwise → sort only by file (natural) and table index
            """
            if tables_order:
                mo = str(entry["sheet_candidate"]).strip()
                mo_pos = mo_rank.get(mo, len(mo_rank) + 1)
                return (mo_pos, file_rank.get(entry["log_file"], 10 ** 9), int(entry["idx_in_file"]))
            else:
                return (file_rank.get(entry["log_file"], 10 ** 9), int(entry["idx_in_file"]), 0)

        table_entries.sort(key=entry_sort_key)

        # =====================================================================
        #                PHASE 3: Assign unique sheet names
        # =====================================================================
        used_sheet_names: set = set(["Summary"])
        for entry in table_entries:
            base_name = self._sanitize_sheet_name(str(entry["sheet_candidate"]))
            final_sheet = self._unique_sheet_name(base_name, used_sheet_names)
            used_sheet_names.add(final_sheet)
            entry["final_sheet"] = final_sheet

        # =====================================================================
        #                PHASE 4: Build the Summary sheet
        # =====================================================================
        summary_rows: List[Dict[str, object]] = []
        for entry in table_entries:
            note = str(entry.get("note", ""))
            separator_str, encoding_str = "", ""

            # Split "Header=..., | encoding=..." into two separate columns
            if note:
                parts = [p.strip() for p in note.split("|")]
                for part in parts:
                    pl = part.lower()
                    if pl.startswith("header=") or "separated" in pl:
                        separator_str = part
                    elif pl.startswith("encoding="):
                        encoding_str = part.replace("encoding=", "")

            df: pd.DataFrame = entry["df"]
            summary_rows.append({
                "File": entry["log_file"],
                "Sheet": entry["final_sheet"],
                "Rows": int(len(df)),
                "Columns": int(df.shape[1]),
                "Separator": separator_str,
                "Encoding": encoding_str,
                "LogFile": entry["log_file"],
                "TablesInLog": entry["tables_in_log"],
            })

        # =====================================================================
        #        PHASE 4.1: Prepare pivot tables for extra summary sheets
        # =====================================================================
        # Collect dataframes for the specific MOs we need
        mo_collectors: Dict[str, List[pd.DataFrame]] = {
            "GUtranSyncSignalFrequency": [],
            "GUtranFreqRelation": [],      # for LTE freq relation checks
            "NRCellDU": [],
            "NRFrequency": [],
            "NRFreqRelation": [],
            "NRSectorCarrier": [],        # for N77B ARFCN checks
            "EndcDistrProfile": [],       # for gUtranFreqRef checks
        }
        for entry in table_entries:
            mo_name = str(entry.get("sheet_candidate", "")).strip()
            if mo_name in mo_collectors:
                df_mo = entry["df"]
                if isinstance(df_mo, pd.DataFrame) and not df_mo.empty:
                    mo_collectors[mo_name].append(df_mo)

        # ---- Build pivots ----
        # Pivot NRCellDU
        df_nr_cell_du = self._concat_or_empty(mo_collectors["NRCellDU"])
        pivot_nr_cells_du = self._safe_pivot_count(
            df=df_nr_cell_du,
            index_field="NodeId",
            columns_field="ssbFrequency",
            values_field="NRCellDUId",
            add_margins=True,
            margins_name="Total",
        )
        pivot_nr_cells_du = self._apply_frequency_column_filter(pivot_nr_cells_du, freq_filters)

        # Pivot NRFrequency
        df_nr_freq = self._concat_or_empty(mo_collectors["NRFrequency"])
        pivot_nr_freq = self._safe_pivot_count(
            df=df_nr_freq,
            index_field="NodeId",
            columns_field="arfcnValueNRDl",
            values_field="NRFrequencyId",
            add_margins=True,
            margins_name="Total",
        )
        pivot_nr_freq = self._apply_frequency_column_filter(pivot_nr_freq, freq_filters)

        # Pivot NRFreqRelation
        df_nr_freq_rel = self._concat_or_empty(mo_collectors["NRFreqRelation"])
        pivot_nr_freq_rel = self._safe_pivot_count(
            df=df_nr_freq_rel,
            index_field="NodeId",
            columns_field="NRFreqRelationId",
            values_field="NRCellCUId",
            add_margins=True,
            margins_name="Total",
        )
        pivot_nr_freq_rel = self._apply_frequency_column_filter(pivot_nr_freq_rel, freq_filters)

        # Pivot GUtranSyncSignalFrequency
        df_gu_sync_signal_freq = self._concat_or_empty(mo_collectors["GUtranSyncSignalFrequency"])
        pivot_gu_sync_signal_freq = self._safe_crosstab_count(
            df=df_gu_sync_signal_freq,
            index_field="NodeId",
            columns_field="arfcn",
            add_margins=True,
            margins_name="Total",
        )
        pivot_gu_sync_signal_freq = self._apply_frequency_column_filter(pivot_gu_sync_signal_freq, freq_filters)

        # Extra tables for audit
        df_gu_freq_rel = self._concat_or_empty(mo_collectors["GUtranFreqRelation"])
        df_nr_sector_carrier = self._concat_or_empty(mo_collectors["NRSectorCarrier"])
        df_endc_distr_profile = self._concat_or_empty(mo_collectors["EndcDistrProfile"])

        # =====================================================================
        #                PHASE 4.2: Build SummaryAudit
        # =====================================================================
        summary_audit_df = self._build_summary_audit(
            df_nr_cell_du=df_nr_cell_du,
            df_nr_freq=df_nr_freq,
            df_nr_freq_rel=df_nr_freq_rel,
            df_gu_sync_signal_freq=df_gu_sync_signal_freq,
            df_gu_freq_rel=df_gu_freq_rel,
            df_nr_sector_carrier=df_nr_sector_carrier,
            df_endc_distr_profile=df_endc_distr_profile,
        )

        # =====================================================================
        #                PHASE 5: Write the Excel file
        # =====================================================================
        with pd.ExcelWriter(excel_path, engine="openpyxl") as writer:
            # Write Summary first
            pd.DataFrame(summary_rows).to_excel(writer, sheet_name="Summary", index=False)

            # Extra summary sheets
            pivot_nr_cells_du.to_excel(writer, sheet_name="Summary NR_CelDU", index=False)
            pivot_nr_freq.to_excel(writer, sheet_name="Summary NR_Frequency", index=False)
            pivot_nr_freq_rel.to_excel(writer, sheet_name="Summary NR_FreqRelation", index=False)
            pivot_gu_sync_signal_freq.to_excel(writer, sheet_name="Summary GU_SyncSignalFrequency", index=False)

            # SummaryAudit with high-level checks
            summary_audit_df.to_excel(writer, sheet_name="SummaryAudit", index=False)

            # Then write each table in the final determined order
            for entry in table_entries:
                entry["df"].to_excel(writer, sheet_name=entry["final_sheet"], index=False)

            # Color the 'Summary*' tabs in green
            color_summary_tabs(writer, prefix="Summary", rgb_hex="00B050")

            # Enable filters (and freeze header row) on all sheets
            enable_header_filters(writer, freeze_header=True)

        print(f"{module_name} Wrote Excel with {len(table_entries)} sheet(s) in: '{excel_path}'")

        # =====================================================================
        #                PHASE 6: Generate PPT textual summary
        # =====================================================================
        try:
            ppt_path = self._generate_ppt_summary(summary_audit_df, excel_path, module_name)
            if ppt_path:
                print(f"{module_name} PPT summary generated in: '{ppt_path}'")
        except Exception as ex:
            # Never fail the whole module just for PPT creation
            print(f"{module_name} [WARN] PPT summary generation failed: {ex}")

        return excel_path

    # =====================================================================
    #                        PRIVATE HELPERS (I/O)
    # =====================================================================
    def _find_log_files(self, folder: str) -> List[str]:
        files = []
        for name in os.listdir(folder):
            lower = name.lower()
            if lower.endswith((".log", ".logs", ".txt")):
                p = os.path.join(folder, name)
                if os.path.isfile(p):
                    files.append(p)
        files.sort()
        return files

    def _read_text_file(self, path: str) -> Tuple[List[str], Optional[str]]:
        return read_text_with_encoding(path)

    # =====================================================================
    #                        PRIVATE HELPERS (Parsing)
    # =====================================================================
    def _parse_log_lines(self, lines: List[str], forced_header_idx: Optional[int] = None) -> Tuple[pd.DataFrame, str]:
        valid = [ln for ln in lines if ln.strip() and not self.SUMMARY_RE.match(ln)]
        header_idx = forced_header_idx
        if header_idx is None:
            header_idx = self._fallback_header_index(valid, lines)
        if header_idx is None:
            return pd.DataFrame(), "No header detected"

        header_line = lines[header_idx].strip()
        any_tab = any("\t" in ln for ln in valid)
        data_sep: Optional[str] = "\t" if any_tab else ("," if any("," in ln for ln in valid) else None)

        if header_line.startswith("SubNetwork"):
            header_cols = [c.strip() for c in header_line.split(",")]
        else:
            header_cols = [c.strip() for c in (header_line.split(data_sep) if data_sep else re.split(r"\s+", header_line.strip()))]
        header_cols = make_unique_columns(header_cols)

        rows: List[List[str]] = []
        for ln in lines[header_idx + 1:]:
            if not ln.strip() or self.SUMMARY_RE.match(ln):
                continue
            parts = [p.strip() for p in (ln.split(data_sep) if data_sep else re.split(r"\s+", ln.strip()))]
            if len(parts) < len(header_cols):
                parts += [""] * (len(header_cols) - len(parts))
            elif len(parts) > len(header_cols):
                parts = parts[:len(header_cols)]
            rows.append(parts)

        df = pd.DataFrame(rows, columns=header_cols)
        df = df.replace({"nan": "", "NaN": "", "None": "", "NULL": ""}).dropna(how="all")
        for c in df.columns:
            df[c] = df[c].astype(str).str.strip()

        note = "Header=SubNetwork-comma" if header_line.startswith("SubNetwork") else (
            "Tab-separated" if data_sep == "\t" else ("Comma-separated" if data_sep == "," else "Whitespace-separated")
        )
        return df, note

    def _fallback_header_index(self, valid_lines: List[str], all_lines: List[str]) -> Optional[int]:
        any_tab = any("\t" in ln for ln in valid_lines)
        sep: Optional[str] = "\t" if any_tab else ("," if any("," in ln for ln in valid_lines) else None)
        for i, ln in enumerate(all_lines):
            if not ln.strip() or self.SUMMARY_RE.match(ln):
                continue
            if sep == "\t" and "\t" in ln:
                return i
            if sep == "," and "," in ln:
                return i
            if sep is None:
                return i
        return None

    @staticmethod
    def _find_subnetwork_header_index(lines: List[str]) -> Optional[int]:
        for i, ln in enumerate(lines):
            if ln.strip().startswith("SubNetwork"):
                return i
        return None

    @staticmethod
    def _extract_mo_name_from_previous_line(lines: List[str], header_idx: Optional[int]) -> Optional[str]:
        if header_idx is None or header_idx == 0:
            return None
        prev = lines[header_idx - 1].strip()
        if not prev:
            return None
        if "," in prev:
            last = prev.split(",")[-1].strip()
            return last or None
        toks = prev.split()
        return toks[-1].strip() if toks else None

    # =====================================================================
    #                        PRIVATE HELPERS (Sheets)
    # =====================================================================
    @staticmethod
    def _sanitize_sheet_name(name: str) -> str:
        return sanitize_sheet_name(name)

    @staticmethod
    def _unique_sheet_name(base: str, used: set) -> str:
        return unique_sheet_name(base, used)

    @staticmethod
    def _cap_rows(df: pd.DataFrame, note: str, max_rows_excel: int = 1_048_576) -> Tuple[pd.DataFrame, str]:
        if len(df) > max_rows_excel:
            df = df.iloc[:max_rows_excel, :].copy()
            note = (note + " | " if note else "") + f"Trimmed to {max_rows_excel} rows"
        return df, note

    @staticmethod
    def _find_all_subnetwork_headers(lines: List[str]) -> List[int]:
        return find_all_subnetwork_headers(lines)

    # =====================================================================
    #                     PRIVATE HELPERS (Pivots & Filters)
    # =====================================================================
    @staticmethod
    def _concat_or_empty(dfs: List[pd.DataFrame]) -> pd.DataFrame:
        """Return a single concatenated DataFrame or an empty one if none; align on common cols if needed."""
        if not dfs:
            return pd.DataFrame()
        try:
            return pd.concat(dfs, ignore_index=True)
        except Exception:
            common_cols = set.intersection(*(set(d.columns) for d in dfs)) if dfs else set()
            if not common_cols:
                return pd.DataFrame()
            dfs_aligned = [d[list(common_cols)].copy() for d in dfs]
            return pd.concat(dfs_aligned, ignore_index=True)

    @staticmethod
    def _resolve_column_case_insensitive(df: pd.DataFrame, candidates: List[str]) -> Optional[str]:
        """
        Resolve a column name by trying several candidates, case-insensitive and ignoring underscores/spaces.
        """
        if df is None or df.empty:
            return None

        def _canon(s: str) -> str:
            return re.sub(r"[\s_]+", "", str(s).strip().lower())

        cols = list(df.columns)
        canon_map = {_canon(c): c for c in cols}
        for cand in candidates:
            key = _canon(cand)
            if key in canon_map:
                return canon_map[key]
        # Fallback: startswith-based match
        for cand in candidates:
            key = _canon(cand)
            for c in cols:
                if _canon(c).startswith(key):
                    return c
        return None

    @staticmethod
    def _parse_int_frequency(value: object) -> Optional[int]:
        """
        Try to parse a frequency/ARFCN value as integer from the leading numeric part
        of the string (before any non-digit chars like '-' or spaces).
        Examples:
          - '653952-30-20-0-1' -> 653952
          - '648672 some text' -> 648672
          - '  647328'         -> 647328
        """
        if value is None:
            return None
        s = str(value).strip()
        if not s:
            return None

        # Extract leading digits only
        m = re.match(r"^(\d+)", s)
        if not m:
            return None

        try:
            return int(m.group(1))
        except Exception:
            return None

    @staticmethod
    def _is_n77_from_string(value: object) -> bool:
        """
        Determine if a cell can be considered N77 based on ARFCN/SSB string:
        NOTE: here we approximate N77 as frequencies whose textual representation starts with '6'.
        """
        if value is None:
            return False
        s = str(value).strip()
        return bool(s) and s[0] == "6"

    def _safe_pivot_count(
            self,
            df: pd.DataFrame,
            index_field: str,
            columns_field: str,
            values_field: str,
            add_margins: bool = True,
            margins_name: str = "Total",
    ) -> pd.DataFrame:
        """
        Robust pivot builder that prevents 'Grouper for ... not 1-dimensional' errors.
        """
        if df is None or not isinstance(df, pd.DataFrame) or df.empty:
            return pd.DataFrame({"Info": ["Table not found or empty"]})

        if isinstance(df.columns, pd.MultiIndex):
            df.columns = ["_".join([str(c).strip() for c in tup if str(c).strip()]) for tup in df.columns]
        if isinstance(df.index, pd.MultiIndex):
            df = df.reset_index()

        work = df.reset_index(drop=True).copy()

        work.columns = pd.Index([str(c).strip() for c in work.columns])
        seen_lower = set()
        unique_cols = []
        for c in work.columns:
            cl = c.lower()
            if cl in seen_lower:
                continue
            seen_lower.add(cl)
            unique_cols.append(c)
        work = work[unique_cols]

        def _resolve(name: str) -> Optional[str]:
            nl = name.lower()
            for c in work.columns:
                if c.lower() == nl or c.lower().startswith(nl + "_"):
                    return c
            return None

        idx_col = _resolve(index_field)
        col_col = _resolve(columns_field)
        val_col = _resolve(values_field)

        if not all([idx_col, col_col, val_col]):
            missing = [n for n, v in [(index_field, idx_col), (columns_field, col_col), (values_field, val_col)] if v is None]
            return pd.DataFrame({
                "Info": [f"Required columns missing: {', '.join(missing)}"],
                "PresentColumns": [", ".join(work.columns.tolist())],
            })

        for col in {idx_col, col_col, val_col}:
            work[col] = work[col].astype(str).str.strip()

        try:
            piv = pd.pivot_table(
                work,
                index=idx_col,
                columns=col_col,
                values=val_col,
                aggfunc="count",
                fill_value=0,
                margins=add_margins,
                margins_name=margins_name,
            ).reset_index()

            if isinstance(piv.columns, pd.MultiIndex):
                piv.columns = [" ".join([str(x) for x in tup if str(x)]).strip() for tup in piv.columns]

            return piv

        except Exception as ex:
            return pd.DataFrame({
                "Error": [f"Pivot build failed: {ex}"],
                "PresentColumns": [", ".join(work.columns.tolist())],
            })

    def _safe_crosstab_count(
            self,
            df: pd.DataFrame,
            index_field: str,
            columns_field: str,
            add_margins: bool = True,
            margins_name: str = "Total",
    ) -> pd.DataFrame:
        """
        Build a frequency table with pd.crosstab (no 'values' needed).
        """
        import unicodedata
        import re as re_local

        if df is None or not isinstance(df, pd.DataFrame) or df.empty:
            return pd.DataFrame({"Info": ["Table not found or empty"]})

        work = df.copy()
        if isinstance(work.columns, pd.MultiIndex):
            work.columns = ["_".join([str(c) for c in tup if str(c)]).strip() for tup in work.columns]
        if isinstance(work.index, pd.MultiIndex):
            work = work.reset_index()
        work = work.reset_index(drop=True)

        def _norm_header(s: str) -> str:
            s = "" if s is None else str(s)
            s = unicodedata.normalize("NFKC", s).replace("\ufeff", "").replace("\u200b", "").replace("\xa0", " ")
            s = re_local.sub(r"\s+", " ", s.strip())
            return s

        work.columns = pd.Index([_norm_header(c) for c in work.columns])

        def _canon(s: str) -> str:
            s = s.lower().replace(" ", "").replace("_", "").replace("-", "")
            return s

        seen = set()
        keep = []
        for c in work.columns:
            k = _canon(c)
            if k in seen:
                continue
            seen.add(k)
            keep.append(c)
        work = work[keep]

        def _resolve(name: str) -> Optional[str]:
            target = _canon(_norm_header(name))
            for c in work.columns:
                if _canon(c) == target:
                    return c
            for c in work.columns:
                if _canon(c).startswith(target):
                    return c
            return None

        idx_col = _resolve(index_field)
        col_col = _resolve(columns_field)
        if not idx_col or not col_col:
            missing = [n for n, v in [(index_field, idx_col), (columns_field, col_col)] if v is None]
            return pd.DataFrame({
                "Info": [f"Required columns missing: {', '.join(missing)}"],
                "PresentColumns": [", ".join(work.columns.tolist())],
            })

        work[idx_col] = work[idx_col].astype(str).map(_norm_header)
        work[col_col] = work[col_col].astype(str).map(_norm_header)

        try:
            ct = pd.crosstab(
                index=work[idx_col],
                columns=work[col_col],
                dropna=False,
            ).reset_index()

            if add_margins:
                ct["Total"] = ct.drop(columns=[idx_col]).sum(axis=1)
                total_row = {idx_col: "Total"}
                for c in ct.columns:
                    if c != idx_col:
                        total_row[c] = ct[c].sum()
                ct = pd.concat([ct, pd.DataFrame([total_row])], ignore_index=True)

            return ct
        except Exception as ex:
            return pd.DataFrame({
                "Error": [f"Crosstab build failed: {ex}"],
                "PresentColumns": [", ".join(work.columns.tolist())],
            })

    @staticmethod
    def _apply_frequency_column_filter(piv: pd.DataFrame, filters: List[str]) -> pd.DataFrame:
        """
        Keep only the first (index) column, 'Total' (if present), and columns whose
        header contains any of the provided substrings (case-insensitive).
        """
        if not isinstance(piv, pd.DataFrame) or piv.empty or not filters:
            return piv

        cols = [str(c) for c in piv.columns.tolist()]
        keep = []

        if cols:
            keep.append(cols[0])

        fl = [f.lower() for f in filters if f]
        for c in cols[1:]:
            lc = c.lower()
            if c == "Total" or lc == "total":
                keep.append(c)
                continue
            if any(f in lc for f in fl):
                keep.append(c)

        if len(keep) <= 1 and "Total" in cols and "Total" not in keep:
            keep.append("Total")

        try:
            return piv.loc[:, keep]
        except Exception:
            return piv

    # =====================================================================
    #                     PRIVATE HELPERS (SummaryAudit)
    # =====================================================================
    def _is_allowed_n77_ssb(self, v: object) -> bool:
        freq = self._parse_int_frequency(v)
        return freq in self.ALLOWED_N77_SSB if freq is not None else False

    def _is_allowed_n77_arfcn(self, v: object) -> bool:
        freq = self._parse_int_frequency(v)
        return freq in self.ALLOWED_N77_ARFCN if freq is not None else False

    def _is_not_old_not_new(self, v: object) -> bool:
        freq = self._parse_int_frequency(v)
        return freq not in (self.OLD_ARFCN, self.NEW_ARFCN)

    def _only_not_old_not_new(self, series):
        return all(self._is_not_old_not_new(v) for v in series)

    def _is_new(self, v: object) -> bool:
        freq = self._parse_int_frequency(v)
        return freq == self.NEW_ARFCN

    def _is_old(self, v: object) -> bool:
        freq = self._parse_int_frequency(v)
        return freq == self.OLD_ARFCN

    def _build_summary_audit(
        self,
        df_nr_cell_du: pd.DataFrame,
        df_nr_freq: pd.DataFrame,
        df_nr_freq_rel: pd.DataFrame,
        df_gu_sync_signal_freq: pd.DataFrame,
        df_gu_freq_rel: pd.DataFrame,
        df_nr_sector_carrier: pd.DataFrame,
        df_endc_distr_profile: pd.DataFrame,
    ) -> pd.DataFrame:
        """
        Build a synthetic 'SummaryAudit' table with high-level checks:

        - Count/list nodes with N77 cells (NRCellDU / NRSectorCarrier)
        - Count NR/LTE nodes where specific ARFCNs (NEW_ARFCN / OLD_ARFCN) are defined
        - Check for references to new SSB/ARFCNs in Frequency/FreqRelation tables
        - Check cardinality limits per cell and per node
        - Check EndcDistrProfile gUtranFreqRef values

        NOTE:
        - N77 cells are approximated as those with ARFCN/SSB text starting with '6'.
        - This function is best-effort and will not raise exceptions; any error is
          represented as a row in the resulting dataframe.
        """
        rows: List[Dict[str, object]] = []

        def add_row(category: str, metric: str, value: object, extra: str = "") -> None:
            rows.append({
                "Category": category,
                "Metric": metric,
                "Value": value,
                "ExtraInfo": extra,
            })

        # ----------------------------- N77 CELLS (NRCellDU / NRSectorCarrier) -----------------------------
        try:
            # NRCellDU: N77 cells by SSB
            if df_nr_cell_du is not None and not df_nr_cell_du.empty:
                node_col_nr = self._resolve_column_case_insensitive(df_nr_cell_du, ["NodeId"])
                ssb_col = self._resolve_column_case_insensitive(df_nr_cell_du, ["ssbFrequency", "ssb", "arfcnSsb", "arfcn"])
                if node_col_nr and ssb_col:
                    work = df_nr_cell_du[[node_col_nr, ssb_col]].copy()
                    work[node_col_nr] = work[node_col_nr].astype(str)
                    work[ssb_col] = work[ssb_col].astype(str)

                    mask_n77 = work[ssb_col].map(self._is_n77_from_string)
                    n77_nodes = sorted(set(work.loc[mask_n77, node_col_nr]))
                    add_row(
                        "NRCellDU",
                        "N77 nodes (NRCellDU SSB starting with '6')",
                        len(n77_nodes),
                        ", ".join(n77_nodes),
                    )

                    # Cells with N77 SSB but not in allowed list (if provided)
                    if self.ALLOWED_N77_SSB:
                        invalid_mask = mask_n77 & ~work[ssb_col].map(self._is_allowed_n77_ssb)
                        invalid_rows = work.loc[invalid_mask, [node_col_nr, ssb_col]]
                        if not invalid_rows.empty:
                            add_row(
                                "NRCellDU",
                                "N77 nodes with SSB not in allowed list",
                                len(invalid_rows),
                                "; ".join(
                                    f"{r[node_col_nr]}: {r[ssb_col]}"
                                    for _, r in invalid_rows.head(50).iterrows()
                                ) + (" (truncated)" if len(invalid_rows) > 50 else ""),
                            )
                else:
                    add_row("NRCellDU", "NRCellDU table present but required columns missing", "N/A")
            else:
                add_row("NRCellDU", "NRCellDU table", "Table not found or empty")
        except Exception as ex:
            add_row("NRCellDU", "Error while checking NRCellDU", f"ERROR: {ex}")

        try:
            # NRSectorCarrier: N77 sectors by ARFCN
            if df_nr_sector_carrier is not None and not df_nr_sector_carrier.empty:
                node_col_sec = self._resolve_column_case_insensitive(df_nr_sector_carrier, ["NodeId"])
                arfcn_col_sec = self._resolve_column_case_insensitive(df_nr_sector_carrier, ["arfcnDL", "arfcn", "arfcnValueNRDl"])
                if node_col_sec and arfcn_col_sec:
                    work = df_nr_sector_carrier[[node_col_sec, arfcn_col_sec]].copy()
                    work[node_col_sec] = work[node_col_sec].astype(str)
                    work[arfcn_col_sec] = work[arfcn_col_sec].astype(str)

                    mask_n77b = work[arfcn_col_sec].map(self._is_n77_from_string)
                    n77b_nodes = sorted(set(work.loc[mask_n77b, node_col_sec]))
                    add_row(
                        "NRSectorCarrier",
                        "N77 nodes with ARFCN starting with '6'",
                        len(n77b_nodes),
                        ", ".join(n77b_nodes),
                    )

                    if self.ALLOWED_N77_ARFCN:
                        invalid_mask = mask_n77b & ~work[arfcn_col_sec].map(self._is_allowed_n77_arfcn)
                        invalid_rows = work.loc[invalid_mask, [node_col_sec, arfcn_col_sec]]
                        if not invalid_rows.empty:
                            add_row(
                                "NRSectorCarrier",
                                "N77 nodes with ARFCN not in allowed list",
                                len(invalid_rows),
                                "; ".join(
                                    f"{r[node_col_sec]}: {r[arfcn_col_sec]}"
                                    for _, r in invalid_rows.head(50).iterrows()
                                ) + (" (truncated)" if len(invalid_rows) > 50 else ""),
                            )
                else:
                    add_row("NRSectorCarrier", "NRSectorCarrier table present but required columns missing", "N/A")
            else:
                add_row("NRSectorCarrier", "NRSectorCarrier table", "Table not found or empty")
        except Exception as ex:
            add_row("NRSectorCarrier", "Error while checking NRSectorCarrier", f"ERROR: {ex}")

        # ----------------------------- NR FREQUENCY (OLD/NEW ARFCN) -----------------------------
        try:
            if df_nr_freq is not None and not df_nr_freq.empty:
                node_col = self._resolve_column_case_insensitive(df_nr_freq, ["NodeId"])
                arfcn_col = self._resolve_column_case_insensitive(df_nr_freq, ["arfcnValueNRDl"])
                if node_col and arfcn_col:
                    work = df_nr_freq[[node_col, arfcn_col]].copy()
                    work[node_col] = work[node_col].astype(str)

                    grouped = work.groupby(node_col)[arfcn_col]
                    mask = grouped.apply(self._only_not_old_not_new)
                    not_old_not_new_nodes = sorted(mask[mask].index.astype(str))
                    new_nodes = sorted(set(work.loc[work[arfcn_col].map(self._is_new), node_col]))
                    old_nodes = sorted(set(work.loc[work[arfcn_col].map(self._is_old), node_col]))

                    add_row(
                        "NRFrequency",
                        f"N77 nodes with the ARFCN not in ({self.OLD_ARFCN}, {self.NEW_ARFCN}) in NRFrequency",
                        len(not_old_not_new_nodes),
                        ", ".join(not_old_not_new_nodes),
                    )
                    add_row(
                        "NRFrequency",
                        f"N77 nodes with the old ARFCN ({self.OLD_ARFCN}) in NRFrequency",
                        len(old_nodes),
                        ", ".join(old_nodes),
                    )
                    add_row(
                        "NRFrequency",
                        f"N77 nodes with the new ARFCN ({self.NEW_ARFCN}) in NRFrequency",
                        len(new_nodes),
                        ", ".join(new_nodes),
                    )
                else:
                    add_row("NRFrequency", "NRFrequency table present but required columns missing", "N/A")
            else:
                add_row("NRFrequency", "NRFrequency table", "Table not found or empty")
        except Exception as ex:
            add_row("NRFrequency", "Error while checking NRFrequency", f"ERROR: {ex}")

        # ----------------------------- NR FREQRELATION (OLD/NEW ARFCN) -----------------------------
        try:
            if df_nr_freq_rel is not None and not df_nr_freq_rel.empty:
                node_col = self._resolve_column_case_insensitive(df_nr_freq_rel, ["NodeId"])
                arfcn_col = self._resolve_column_case_insensitive(df_nr_freq_rel, ["NRFreqRelationId"])
                if node_col and arfcn_col:
                    work = df_nr_freq_rel[[node_col, arfcn_col]].copy()
                    work[node_col] = work[node_col].astype(str)

                    grouped = work.groupby(node_col)[arfcn_col]
                    mask = grouped.apply(self._only_not_old_not_new)
                    not_old_not_new_nodes = sorted(mask[mask].index.astype(str))
                    new_nodes = sorted(set(work.loc[work[arfcn_col].map(self._is_new), node_col]))
                    old_nodes = sorted(set(work.loc[work[arfcn_col].map(self._is_old), node_col]))

                    add_row(
                        "NRFreqRelation",
                        f"N77 nodes with the ARFCN not in ({self.OLD_ARFCN}, {self.NEW_ARFCN}) in NRFreqRelation",
                        len(not_old_not_new_nodes),
                        ", ".join(not_old_not_new_nodes),
                    )
                    add_row(
                        "NRFreqRelation",
                        f"NR nodes with the old ARFCN ({self.OLD_ARFCN}) in NRFreqRelation",
                        len(old_nodes),
                        ", ".join(old_nodes),
                    )
                    add_row(
                        "NRFreqRelation",
                        f"NR nodes with the old ARFCN ({self.NEW_ARFCN}) in NRFreqRelation",
                        len(new_nodes),
                        ", ".join(new_nodes),
                    )

                else:
                    add_row("NRFreqRelation", "NRFreqRelation table present but ARFCN column missing", "N/A")
            else:
                add_row("NRFreqRelation", "NRFreqRelation table", "Table not found or empty")
        except Exception as ex:
            add_row("NRFreqRelation", "Error while checking NRFreqRelation", f"ERROR: {ex}")

        # ----------------------------- LTE GUtranSyncSignalFrequency -----------------------------
        try:
            if df_gu_sync_signal_freq is not None and not df_gu_sync_signal_freq.empty:
                node_col = self._resolve_column_case_insensitive(df_gu_sync_signal_freq, ["NodeId"])
                arfcn_col = self._resolve_column_case_insensitive(df_gu_sync_signal_freq, ["arfcn", "arfcnDL"])
                if node_col and arfcn_col:
                    work = df_gu_sync_signal_freq[[node_col, arfcn_col]].copy()
                    work[node_col] = work[node_col].astype(str)

                    grouped = work.groupby(node_col)[arfcn_col]
                    mask = grouped.apply(self._only_not_old_not_new)
                    not_old_not_new_nodes = sorted(mask[mask].index.astype(str))
                    new_nodes = sorted(set(work.loc[work[arfcn_col].map(self._is_new), node_col]))
                    old_nodes = sorted(set(work.loc[work[arfcn_col].map(self._is_old), node_col]))

                    add_row(
                        "GUtranSyncSignalFrequency",
                        f"LTE nodes with the ARFCN not in ({self.OLD_ARFCN}, {self.NEW_ARFCN}) in GUtranSyncSignalFrequency",
                        len(not_old_not_new_nodes),
                        ", ".join(not_old_not_new_nodes),
                    )
                    add_row(
                        "GUtranSyncSignalFrequency",
                        f"LTE nodes with the old ARFCN ({self.OLD_ARFCN}) in GUtranSyncSignalFrequency",
                        len(old_nodes),
                        ", ".join(old_nodes),
                    )
                    add_row(
                        "GUtranSyncSignalFrequency",
                        f"LTE nodes with the new ARFCN ({self.NEW_ARFCN}) in GUtranSyncSignalFrequency",
                        len(new_nodes),
                        ", ".join(new_nodes),
                    )
                else:
                    add_row("GUtranSyncSignalFrequency", "GUtranSyncSignalFrequency table present but required columns missing", "N/A")
            else:
                add_row("GUtranSyncSignalFrequency", "GUtranSyncSignalFrequency table", "Table not found or empty")
        except Exception as ex:
            add_row("GUtranSyncSignalFrequency", "Error while checking GUtranSyncSignalFrequency", f"ERROR: {ex}")

        # ----------------------------- LTE GUtranFreqRelation -----------------------------
        try:
            if df_gu_freq_rel is not None and not df_gu_freq_rel.empty:
                node_col = self._resolve_column_case_insensitive(df_gu_freq_rel, ["NodeId"])
                arfcn_col = self._resolve_column_case_insensitive(df_gu_freq_rel, ["GUtranFreqRelationId", "gUtranFreqRelationId"])
                if node_col and arfcn_col:
                    work = df_gu_freq_rel[[node_col, arfcn_col]].copy()
                    work[node_col] = work[node_col].astype(str)

                    grouped = work.groupby(node_col)[arfcn_col]
                    mask = grouped.apply(self._only_not_old_not_new)
                    not_old_not_new_nodes = sorted(mask[mask].index.astype(str))
                    new_nodes = sorted(set(work.loc[work[arfcn_col].map(self._is_new), node_col]))
                    old_nodes = sorted(set(work.loc[work[arfcn_col].map(self._is_old), node_col]))

                    add_row(
                        "GUtranFreqRelation",
                        f"LTE nodes with the ARFCN not in ({self.OLD_ARFCN}, {self.NEW_ARFCN}) in GUtranFreqRelation",
                        len(not_old_not_new_nodes),
                        ", ".join(not_old_not_new_nodes),
                    )
                    add_row(
                        "GUtranFreqRelation",
                        f"LTE nodes with the old ARFCN ({self.OLD_ARFCN}) in GUtranFreqRelation",
                        len(old_nodes),
                        ", ".join(old_nodes),
                    )
                    add_row(
                        "GUtranFreqRelation",
                        f"LTE nodes with the new ARFCN ({self.NEW_ARFCN}) in GUtranFreqRelation",
                        len(new_nodes),
                        ", ".join(new_nodes),
                    )

                else:
                    add_row("GUtranFreqRelation", "GUtranFreqRelation table present but ARFCN/NodeId missing", "N/A")
            else:
                add_row("GUtranFreqRelation", "GUtranFreqRelation table", "Table not found or empty")
        except Exception as ex:
            add_row("GUtranFreqRelation", "Error while checking GUtranFreqRelation", f"ERROR: {ex}")

        # ----------------------------- CARDINALITY LIMITS -----------------------------
        # Max 16 NRFreqRelation and 16 GUtranFreqRelation per cell
        try:
            if df_nr_freq_rel is not None and not df_nr_freq_rel.empty:
                cell_col = self._resolve_column_case_insensitive(df_nr_freq_rel, ["NRCellCUId", "NRCellId", "CellId"])
                if cell_col:
                    counts = df_nr_freq_rel[cell_col].value_counts(dropna=False)
                    max_count = int(counts.max()) if not counts.empty else 0
                    over_limit = counts[counts >= 16]
                    add_row(
                        "Cardinality",
                        "Max NRFreqRelation per NR cell (limit 16)",
                        max_count,
                        "; ".join(f"{idx}: {cnt}" for idx, cnt in over_limit.head(50).items())
                        + (" (truncated)" if len(over_limit) > 50 else ""),
                    )
                else:
                    add_row("Cardinality", "NRFreqRelation per cell (required cell column missing)", "N/A")
            else:
                add_row("Cardinality", "NRFreqRelation per cell", "Table not found or empty")
        except Exception as ex:
            add_row("Cardinality", "Error while checking NRFreqRelation cardinality", f"ERROR: {ex}")

        try:
            if df_gu_freq_rel is not None and not df_gu_freq_rel.empty:
                cell_col_gu = self._resolve_column_case_insensitive(df_gu_freq_rel, ["EUtranCellFDDId", "EUtranCellId", "CellId", "GUCellId"])
                if cell_col_gu:
                    counts = df_gu_freq_rel[cell_col_gu].value_counts(dropna=False)
                    max_count = int(counts.max()) if not counts.empty else 0
                    over_limit = counts[counts >= 16]
                    add_row(
                        "Cardinality",
                        "Max GUtranFreqRelation per LTE cell (limit 16)",
                        max_count,
                        "; ".join(f"{idx}: {cnt}" for idx, cnt in over_limit.head(50).items())
                        + (" (truncated)" if len(over_limit) > 50 else ""),
                    )
                else:
                    add_row("Cardinality", "GUtranFreqRelation per LTE cell (required cell column missing)", "N/A")
            else:
                add_row("Cardinality", "GUtranFreqRelation per LTE cell", "Table not found or empty")
        except Exception as ex:
            add_row("Cardinality", "Error while checking GUtranFreqRelation cardinality", f"ERROR: {ex}")

        # Max 24 GUtranSyncSignalFrequency per node
        try:
            if df_gu_sync_signal_freq is not None and not df_gu_sync_signal_freq.empty:
                node_col = self._resolve_column_case_insensitive(df_gu_sync_signal_freq, ["NodeId"])
                if node_col:
                    counts = df_gu_sync_signal_freq[node_col].astype(str).value_counts(dropna=False)
                    max_count = int(counts.max()) if not counts.empty else 0
                    over_limit_nodes = counts[counts >= 24]
                    add_row(
                        "Cardinality",
                        "Max GUtranSyncSignalFrequency definitions per node (limit 24)",
                        max_count,
                        "; ".join(f"{idx}: {cnt}" for idx, cnt in over_limit_nodes.head(50).items())
                        + (" (truncated)" if len(over_limit_nodes) > 50 else ""),
                    )
                else:
                    add_row("Cardinality", "GUtranSyncSignalFrequency per node (NodeId missing)", "N/A")
            else:
                add_row("Cardinality", "GUtranSyncSignalFrequency per node", "Table not found or empty")
        except Exception as ex:
            add_row("Cardinality", "Error while checking GUtranSyncSignalFrequency cardinality", f"ERROR: {ex}")

        # Max 64 NRFrequency per node
        try:
            if df_nr_freq is not None and not df_nr_freq.empty:
                node_col = self._resolve_column_case_insensitive(df_nr_freq, ["NodeId"])
                if node_col:
                    counts = df_nr_freq[node_col].astype(str).value_counts(dropna=False)
                    max_count = int(counts.max()) if not counts.empty else 0
                    over_limit_nodes = counts[counts >= 64]
                    add_row(
                        "Cardinality",
                        "Max NRFrequency definitions per node (limit 64)",
                        max_count,
                        "; ".join(f"{idx}: {cnt}" for idx, cnt in over_limit_nodes.head(50).items())
                        + (" (truncated)" if len(over_limit_nodes) > 50 else ""),
                    )
                else:
                    add_row("Cardinality", "NRFrequency per node (NodeId missing)", "N/A")
            else:
                add_row("Cardinality", "NRFrequency per node", "Table not found or empty")
        except Exception as ex:
            add_row("Cardinality", "Error while checking NRFrequency cardinality", f"ERROR: {ex}")

        # ----------------------------- EndcDistrProfile gUtranFreqRef -----------------------------
        try:
            if df_endc_distr_profile is not None and not df_endc_distr_profile.empty:
                node_col_edp = self._resolve_column_case_insensitive(df_endc_distr_profile, ["NodeId"])
                ref_col = self._resolve_column_case_insensitive(df_endc_distr_profile, ["gUtranFreqRef"])
                if node_col_edp and ref_col:
                    work = df_endc_distr_profile[[node_col_edp, ref_col]].copy()
                    work[node_col_edp] = work[node_col_edp].astype(str)
                    work[ref_col] = work[ref_col].astype(str)

                    def _normalize_ref(s: str) -> str:
                        return str(s).replace(" ", "").strip()

                    # Expected pattern is always NEW&other (or equivalent with comma/dash)
                    # We only check that the normalized string contains the numeric NEW_ARFCN at least once.
                    expected_str = str(self.NEW_ARFCN)

                    bad_mask = ~work[ref_col].map(lambda v: expected_str in _normalize_ref(v))
                    bad_rows = work.loc[bad_mask, [node_col_edp, ref_col]]

                    add_row(
                        "EndcDistrProfile",
                        f"EndcDistrProfile rows with gUtranFreqRef not containing {self.NEW_ARFCN}",
                        len(bad_rows),
                        "; ".join(
                            f"{r[node_col_edp]}: {r[ref_col]}"
                            for _, r in bad_rows.head(50).iterrows()
                        ) + (" (truncated)" if len(bad_rows) > 50 else ""),
                    )
                else:
                    add_row("EndcDistrProfile", "EndcDistrProfile table present but NodeId/gUtranFreqRef missing", "N/A")
            else:
                add_row("EndcDistrProfile", "EndcDistrProfile table", "Table not found or empty")
        except Exception as ex:
            add_row("EndcDistrProfile", "Error while checking EndcDistrProfile gUtranFreqRef", f"ERROR: {ex}")

        # If nothing was added, return at least an informational row
        if not rows:
            rows.append({
                "Category": "Info",
                "Metric": "SummaryAudit",
                "Value": "No data available",
                "ExtraInfo": "",
            })

        return pd.DataFrame(rows)

    # =====================================================================
    #                     PRIVATE HELPERS (Summary → PPT)
    # =====================================================================
    @staticmethod
    def _build_text_summary_structure(summary_audit_df: pd.DataFrame) -> Dict[str, List[str]]:
        """
        Transform SummaryAudit rows into a dict structure suitable for PPT generation.

        Returns:
          {
            "Category1": [
                "Metric1: value | extra",
                "Metric2: value | extra",
                ...
            ],
            "Category2": [
                ...
            ],
          }
        """
        sections: Dict[str, List[str]] = {}

        if summary_audit_df is None or summary_audit_df.empty:
            sections["Info"] = ["No audit data available to build textual summary"]
            return sections

        for _, row in summary_audit_df.iterrows():
            category = str(row.get("Category", "") or "Info")
            metric = str(row.get("Metric", "") or "")
            value = row.get("Value", "")
            extra = str(row.get("ExtraInfo", "") or "")

            base_line = f"{metric}: {value}"
            if extra:
                base_line = f"{base_line} | {extra}"

            sections.setdefault(category, []).append(base_line)

        return sections

    def _generate_ppt_summary(
            self,
            summary_audit_df: pd.DataFrame,
            excel_path: str,
            module_name: str = "",
    ) -> Optional[str]:
        """
        Generate a PPTX file next to the Excel with slides per Category from SummaryAudit.

        - First slide: global title.
        - Then, for each Category:
            • One or more slides.
            • Slide title = Category
            • Body:
                - Main bullet (font size 14 pt): "Metric: Value"
                - Secondary bullets (font size 10 pt): one per cell/node/sector
                  that appears after the '|' separator.
                - Secondary bullets are arranged in two columns:
                    · Max 25 items per column (50 per slide).
                    · If there are more than 50 items, a new slide is created
                      for the remaining items (repeating the same main bullet).
        """
        # Late import to avoid hard dependency if pptx is not installed
        try:
            from pptx import Presentation
            from pptx.util import Pt, Inches
        except ImportError:
            print(f"{module_name} [INFO] python-pptx is not installed. Skipping PPT summary.")
            return None

        # Font sizes for different bullet levels
        MAIN_BULLET_SIZE = Pt(14)
        SUB_BULLET_SIZE = Pt(10)

        sections = self._build_text_summary_structure(summary_audit_df)

        # Derive PPT path from Excel path
        base, _ = os.path.splitext(excel_path)
        ppt_path = base + "_Summary.pptx"

        # Helper: set font size for all runs in a paragraph
        def _set_paragraph_font_size(paragraph, size: Pt) -> None:
            for run in paragraph.runs:
                run.font.size = size

        # Helper: chunk a list into pieces of maximum 'size' elements
        def _chunk_list(items: List[str], size: int) -> List[List[str]]:
            return [items[i:i + size] for i in range(0, len(items), size)]

        # Load user template instead of default blank PPT
        template_path = get_resource_path("ppt_templates/ConfigurationAuditTemplate.pptx")
        try:
            prs = Presentation(template_path)
            print(f"{module_name} Using PPT template: {template_path}")
        except Exception as e:
            print(f"{module_name} [WARN] Could not load PPT template, using default. ({e})")
            prs = Presentation()

        # Find appropriate layouts inside the template
        # You may adjust indices depending on your template
        try:
            title_slide_layout = prs.slide_layouts[0]  # Title layout of your template
            content_layout = prs.slide_layouts[2]  # Title + content layout of your template
        except:
            # Fallback to default mapping
            title_slide_layout = prs.slide_layouts[0]
            content_layout = prs.slide_layouts[1]

        # --- Title slide ---
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None

        title.text = "Configuration Audit Summary"
        if subtitle is not None:
            subtitle.text = os.path.basename(excel_path)

        # --- Category slides ---
        for category, lines in sections.items():
            # If there are no lines for this category, create a single simple slide
            if not lines:
                slide = prs.slides.add_slide(content_layout)
                title_shape = slide.shapes.title
                body = slide.placeholders[1] if len(slide.placeholders) > 1 else None

                title_shape.text = category

                if body is None:
                    continue

                tf = body.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = "No data available for this category."
                p.level = 0
                _set_paragraph_font_size(p, MAIN_BULLET_SIZE)
                continue

            # For each logical line: "Metric: Value | Node1, Node2, Node3, ..."
            # we may create multiple slides if there are many nodes.
            for line in lines:
                raw_text = line or ""

                # Split main text and node list by '|'
                if "|" in raw_text:
                    main_text, extra_text = raw_text.split("|", 1)
                    main_text = main_text.strip()
                    extra_text = extra_text.strip()
                else:
                    main_text = raw_text.strip()
                    extra_text = ""

                # If there is no extra_text, just one simple slide with a main bullet
                if not extra_text:
                    slide = prs.slides.add_slide(content_layout)
                    title_shape = slide.shapes.title
                    body = slide.placeholders[1] if len(slide.placeholders) > 1 else None

                    title_shape.text = category

                    if body is None:
                        continue

                    tf = body.text_frame
                    tf.clear()
                    p_main = tf.paragraphs[0]
                    p_main.text = main_text if main_text else "No data available for this category."
                    p_main.level = 0
                    _set_paragraph_font_size(p_main, MAIN_BULLET_SIZE)
                    continue

                # There is extra_text: build the list of secondary items (nodes/cells/sectors)
                cleaned_extra = extra_text.replace(";", ",")
                items = [t.strip() for t in cleaned_extra.split(",") if t.strip()]

                if not items:
                    # Fallback: treat as no extra_text
                    slide = prs.slides.add_slide(content_layout)
                    title_shape = slide.shapes.title
                    body = slide.placeholders[1] if len(slide.placeholders) > 1 else None

                    title_shape.text = category

                    if body is None:
                        continue

                    tf = body.text_frame
                    tf.clear()
                    p_main = tf.paragraphs[0]
                    p_main.text = main_text if main_text else "No data available for this category."
                    p_main.level = 0
                    _set_paragraph_font_size(p_main, MAIN_BULLET_SIZE)
                    continue

                # Now we have a non-empty list of secondary items.
                # We will:
                #   - Split items into chunks of 50 (max per slide).
                #   - For each chunk:
                #       · Use main placeholder for the main bullet.
                #       · Create two columns of up to 25 items each (secondary bullets).
                chunks_of_50 = _chunk_list(items, 50)

                for chunk_index, chunk_items in enumerate(chunks_of_50):
                    slide = prs.slides.add_slide(content_layout)
                    title_shape = slide.shapes.title
                    body = slide.placeholders[1] if len(slide.placeholders) > 1 else None

                    title_shape.text = category

                    if body is None:
                        continue

                    # Main bullet in placeholder
                    tf_main = body.text_frame
                    tf_main.clear()
                    p_main = tf_main.paragraphs[0]
                    p_main.text = main_text if main_text else "No data available for this category."
                    p_main.level = 0
                    _set_paragraph_font_size(p_main, MAIN_BULLET_SIZE)

                    # Split secondary items into two columns: 25 max each
                    col1_items = chunk_items[:25]
                    col2_items = chunk_items[25:]

                    # Coordinates for the two columns
                    left_margin = Inches(0.8)
                    top = Inches(2.0)
                    col_width = Inches(4.0)
                    col_height = Inches(4.0)
                    gap_between_cols = Inches(0.3)

                    # Left column
                    if col1_items:
                        tx_box1 = slide.shapes.add_textbox(
                            left_margin,
                            top,
                            col_width,
                            col_height,
                        )
                        tf1 = tx_box1.text_frame
                        tf1.clear()

                        first = True
                        for item in col1_items:
                            bullet_text = f"- {item}"
                            if first:
                                p = tf1.paragraphs[0]
                                p.text = bullet_text
                                first = False
                            else:
                                p = tf1.add_paragraph()
                                p.text = bullet_text

                            p.level = 1  # secondary bullet
                            _set_paragraph_font_size(p, SUB_BULLET_SIZE)

                    # Right column
                    if col2_items:
                        tx_box2 = slide.shapes.add_textbox(
                            left_margin + col_width + gap_between_cols,
                            top,
                            col_width,
                            col_height,
                        )
                        tf2 = tx_box2.text_frame
                        tf2.clear()

                        first = True
                        for item in col2_items:
                            bullet_text = f"- {item}"
                            if first:
                                p = tf2.paragraphs[0]
                                p.text = bullet_text
                                first = False
                            else:
                                p = tf2.add_paragraph()
                                p.text = bullet_text

                            p.level = 1  # secondary bullet
                            _set_paragraph_font_size(p, SUB_BULLET_SIZE)

        prs.save(ppt_path)
        return ppt_path


# --------- kept local to preserve current behavior (module-level helper) ----
def make_unique_columns(cols: List[str]) -> List[str]:
    """
    Ensure column names are unique by appending a numeric suffix when needed.
    """
    seen: Dict[str, int] = {}
    unique = []
    for c in cols:
        base = c or "Col"
        if base not in seen:
            seen[base] = 0
            unique.append(base)
        else:
            seen[base] += 1
            unique.append(f"{base}_{seen[base]}")
    return unique
